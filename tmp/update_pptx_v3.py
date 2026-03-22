"""
MBP-T05 - PPT 수정 스크립트 v3 (정밀 수정)

변경점:
- Slide 3 (문제정의): 내용 대폭 보강
- Slide 5 (기술스택): 왼쪽 대형 영역에 state diagram 하나만 삽입
- Slide 6 (구현결과): 기존 4개 이미지 영역 shape 제거 후 플로우차트/시퀀스 이미지 삽입
- Slide 7 (테스트검증): 기존 그리드(행×열) shape 텍스트를 실제 TC 데이터로 교체,
                       오른쪽 성과지표 영역도 업데이트
- Slide 8 (학습교훈): 텍스트 내용 대폭 보강
- Slide 9 (향후계획): 로드맵 단계별, 실무 적용, 아이디어 상세화
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

BASE = r'c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock'
DIAG = os.path.join(BASE, 'docs', 'diagram')
SRC  = os.path.join(BASE, 'docs', 'MBP-T05_ChildLock_결과발표_Updated.pptx')
OUT  = os.path.join(BASE, 'docs', 'MBP-T05_ChildLock_결과발표_Final.pptx')

BLUE_DARK = RGBColor(0x00, 0x35, 0x7e)
BLUE_MID  = RGBColor(0x00, 0x5b, 0xb5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1a, 0x1a, 0x2e)
GREEN     = RGBColor(0x00, 0x96, 0x60)
ASIL_B    = RGBColor(0xd4, 0x45, 0x00)
PASS_GR   = RGBColor(0x00, 0x70, 0x48)
GRAY_LT   = RGBColor(0xf0, 0xf4, 0xf8)


def set_para_text(tf, idx, text, size=None, bold=None, color=None):
    """텍스트 프레임의 특정 문단 텍스트 교체 (첫 번째 run 기준)"""
    if idx >= len(tf.paragraphs):
        return
    para = tf.paragraphs[idx]
    # 기존 run 전부 지우고 새 run 추가
    para.clear()
    run = para.add_run()
    run.text = text
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = 'Malgun Gothic'


def replace_shape_text(shape, new_text, size=None, bold=None, color=None):
    """shape의 텍스트 프레임 전체를 새 텍스트로 교체"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    # 첫 단락 텍스트 교체
    if tf.paragraphs:
        tf.paragraphs[0].clear()
        run = tf.paragraphs[0].add_run()
        run.text = new_text
        run.font.name = 'Malgun Gothic'
        if size:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color
    # 나머지 단락 제거
    from pptx.oxml.ns import qn
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def replace_shape_multiline(shape, lines):
    """shape 텍스트를 여러 줄로 교체. lines = [(text, size, bold, color), ...]"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    # 모든 기존 단락 제거
    for p in list(tf.paragraphs):
        try:
            p._p.getparent().remove(p._p)
        except Exception:
            pass
    # 새 단락 추가
    txBody = tf._txBody
    from pptx.oxml import parse_xml
    for i, (text, size, bold, color) in enumerate(lines):
        from lxml import etree
        from pptx.oxml.ns import nsmap
        p_xml = '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:r><a:t></a:t></a:r></a:p>'
        p_elem = parse_xml(p_xml)
        txBody.append(p_elem)
        # 마지막 추가된 단락
        para = tf.paragraphs[-1]
        run = para.runs[0]
        run.text = text
        run.font.name = 'Malgun Gothic'
        if size:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_picture_fit(slide, img_path, left_in, top_in, max_w_in, max_h_in):
    """PIL로 비율 계산 후 이미지 삽입"""
    try:
        from PIL import Image
        img = Image.open(img_path)
        iw, ih = img.size
        ratio = min(max_w_in / iw, max_h_in / ih)
        nw, nh = iw * ratio, ih * ratio
        cx = left_in + (max_w_in - nw) / 2
        cy = top_in + (max_h_in - nh) / 2
        pic = slide.shapes.add_picture(img_path, Inches(cx), Inches(cy),
                                        Inches(nw), Inches(nh))
        print(f"    Added image: {os.path.basename(img_path)} ({nw:.2f}\" x {nh:.2f}\")")
        return pic
    except Exception as e:
        # PIL 없으면 width에 맞춤
        try:
            pic = slide.shapes.add_picture(img_path,
                    Inches(left_in), Inches(top_in), Inches(max_w_in))
            print(f"    Added image (no PIL): {os.path.basename(img_path)}")
            return pic
        except Exception as e2:
            print(f"    WARN: Could not add {img_path}: {e2}")
            return None


def remove_shape(slide, shape):
    try:
        shape.element.getparent().remove(shape.element)
    except Exception as e:
        print(f"    WARN: remove_shape failed: {e}")


# ─── SLIDE 3: 문제 정의 ──────────────────────────────────────────────────────

def fix_slide3_problem(slide):
    """AS-IS / TO-BE 내용 대폭 보강"""
    print("  [Slide 3] 문제 정의 내용 보강")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        # AS-IS 내용 교체
        if '문제점 1' in t or ('기계식' in t and '운전자' in t):
            replace_shape_multiline(shape, [
                ('▪ 문제점 1: 기계식 차일드 락의 수동 제어 불편', 11, True, DARK),
                ('  운전자가 뒷좌석에 직접 접근하여 개별 조작 필요 → 운전 중 불가', 10, False, DARK),
                ('▪ 문제점 2: 사고·충돌 발생 시 뒷좌석 탈출 지연', 11, True, DARK),
                ('  기계식 잠금 시 비상 해제 수단 없어 인명 피해 위험 가중', 10, False, DARK),
                ('▪ 문제점 3: 후방 위험 접근 정보와 하차 안전 연동 부재', 11, True, DARK),
                ('  후방 이륜차·차량 접근 중에도 뒷좌석 문이 열릴 수 있어 사고 위험', 10, False, DARK),
                ('▪ 문제점 4: 점화 OFF 시 차일드 락 상태 인지 부재', 11, True, DARK),
                ('  하차 직전 잠금 상태 미인지로 아이가 내부에서 열기 불가', 10, False, DARK),
            ])

        # TO-BE 내용 교체
        elif '목표 1' in t or ('전자 제어' in t and '편의성' in t):
            replace_shape_multiline(shape, [
                ('✔ 목표 1: 전자식 통합 차일드 락 제어 (FG-01)', 11, True, BLUE_MID),
                ('  운전석에서 원버튼 ON/OFF + 속도 기반 자동 잠금 (≥15 km/h)', 10, False, DARK),
                ('✔ 목표 2: 사고 시 즉각 자동 해제 (ASIL B, FG-01)', 11, True, BLUE_MID),
                ('  충돌·에어백 신호 동시 감지 시 비상 해제 최우선 수행', 10, False, DARK),
                ('✔ 목표 3: 후방 위험 연동 자동 보호 (FG-02)', 11, True, BLUE_MID),
                ('  후방 레이더 연동 — 위험 감지 시 자동 잠금 유지 & 운전자 경고', 10, False, DARK),
                ('✔ 목표 4: 시동 OFF / 출발 시 상태 알림 (FG-02)', 11, True, BLUE_MID),
                ('  점화 OFF 및 출발 시 뒷좌석 점유 + CL 상태 1회 운전자 알림', 10, False, DARK),
            ])

        # 핵심 연구 질문
        elif '핵심 연구 질문' in t:
            pass  # 제목은 그대로
        elif 'ISO 26262 기능 안전을 준수하면서' in t:
            replace_shape_text(shape,
                'ISO 26262 기능 안전을 준수하면서, 승객의 편의성과 하차 안전을 모두 충족하는\n'
                '지능형 차일드 락 제어 시스템을 어떻게 설계·구현·검증할 것인가?\n'
                '→ V-모델 기반 SRS↔SDD↔Code↔Test 양방향 추적성으로 답한다.',
                size=13, bold=False, color=DARK)


# ─── SLIDE 5: 기술 스택 — 왼쪽 영역에 state diagram만 ──────────────────────

def fix_slide5_arch(slide):
    """왼쪽 큰 RoundedRect(0.5, 1.3, 8.5×5.5) 내부에 state diagram 하나만 삽입"""
    print("  [Slide 5] 기술스택 — 왼쪽에 state diagram 삽입")

    # 기존에 삽입된 이미지 picture들 제거 (이전 실행 결과물)
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            to_remove.append(shape)
    for s in to_remove:
        remove_shape(slide, s)

    # 왼쪽 영역 title 텍스트박스 추가
    txb = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.35), Inches(8.2), Inches(0.38))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '시스템 상태 다이어그램 (State Diagram)'
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = BLUE_DARK
    run.font.name = 'Malgun Gothic'

    # state diagram 이미지 삽입 — 왼쪽 rounded rect 영역(0.5~9.0, 1.3~6.8) 내에 여백 두고 배치
    state_img = os.path.join(DIAG, 'state', 'childlock_state.png')
    add_picture_fit(slide, state_img,
                    left_in=0.7, top_in=1.82,
                    max_w_in=8.1, max_h_in=4.8)


# ─── SLIDE 6: 구현 결과 — 기존 shape 정리 후 이미지 삽입 ────────────────────

def fix_slide6_impl(slide):
    """기존 이미지·텍스트 제거 후 4개 다이어그램 이미지를 4개 RoundedRect 위에 배치"""
    print("  [Slide 6] 구현결과 이미지 삽입")

    # 기존 picture 및 설명 텍스트box 제거
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            to_remove.append(shape)
        elif shape.has_text_frame:
            t = shape.text.strip()
            if '안전 상태' in t or 'AI 기반' in t or '스크린샷' in t:
                to_remove.append(shape)
    for s in to_remove:
        remove_shape(slide, s)

    # 기존 4개 RoundedRect 위치:
    # [3] (0.5, 1.3, 5.8×2.5)  [4] (6.6, 1.3, 5.8×2.5)
    # [5] (0.5, 4.0, 5.8×2.5)  [6] (6.6, 4.0, 5.8×2.5)
    images = [
        (os.path.join(DIAG, 'flow_chart', 'FG-01_door_control.png'),
         0.5, 1.3, 5.8, 2.5),
        (os.path.join(DIAG, 'flow_chart', 'FG-02_exit_safety.png'),
         6.6, 1.3, 5.8, 2.5),
        (os.path.join(DIAG, 'sequence', 'sequence_uc1.png'),
         0.5, 4.0, 5.8, 2.5),
        (os.path.join(DIAG, 'sequence', 'sequence_uc3.png'),
         6.6, 4.0, 5.8, 2.5),
    ]

    # 캡션 텍스트
    captions = [
        ('FG-01 플로우차트: 도어 잠금/해제 제어 (UC-1,2,3,4)', 0.5, 1.3),
        ('FG-02 플로우차트: 하차 안전 보조 (UC-5,6,7)', 6.6, 1.3),
        ('UC-1 시퀀스: 운전자 스위치 ON/OFF 제어 흐름', 0.5, 4.0),
        ('UC-3 시퀀스: 충돌 감지 시 ASIL B 비상 해제 경로', 6.6, 4.0),
    ]

    for (img_path, lx, ty, mw, mh), (cap_text, cx, cy) in zip(images, captions):
        # 캡션 텍스트박스 (RoundedRect 가장 위쪽에 얹음)
        txb = slide.shapes.add_textbox(
            Inches(cx + 0.1), Inches(ty + 0.05), Inches(mw - 0.2), Inches(0.3))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = cap_text
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = BLUE_DARK
        run.font.name = 'Malgun Gothic'
        # 이미지 (캡션 아래 공간에 배치)
        add_picture_fit(slide, img_path, lx + 0.1, ty + 0.35, mw - 0.2, mh - 0.4)


# ─── SLIDE 7: 테스트 & 검증 — 기존 그리드에 TC 데이터 채우기 ─────────────────

def fix_slide7_test(slide):
    """
    기존 슬라이드 구조:
      - 4열 그리드: col0(x=0.7, w=2.5) col1(x=3.2, w=1.8) col2(x=5.0, w=1.8) col3(x=6.8, w=1.2)
      - shape [5~12]: 헤더 행 (rect + textbox 쌍 × 4)
      - shape [13~52]: 데이터 행 5쌍 × 4열 × 5rows = 각 rect+textbox
      - shape [53~]: 오른쪽 성과지표 영역
    헤더와 데이터 shape의 텍스트박스를 실제 TC 데이터로 교체한다.
    """
    print("  [Slide 7] 테스트 케이스 표 내용 교체")

    # 새 헤더 텍스트
    new_headers = ['TC-ID', 'UC', '테스트 항목 / 기법', '결과']

    # TC 데이터 (15건, 5행×4열 x 3묶음 = 15행 in groups?)
    # 실제로는 5행씩 표시되는 구조
    # 기존 shape에는 5개 행이 있음 (shape 13~52: 5행×4열=20쌍의 rect+textbox)
    # 한 행이 4개 col로 이루어진 rect+textbox 쌍 (총 2개씩 = 8 shapes per row)
    # 5행 = 40 shapes (index 13~52)
    # BUT 실제로 15 TC를 한 슬라이드에 다 보여주기엔 공간이 부족하므로
    # 핵심 8건만 (ASIL B 포함) 표시하고, 나머지는 "기타 7건 PASS" 마지막 row에 표기

    # col 0 (w=2.5): Test ID + 테스트 기법
    # col 1 (w=1.8): UC
    # col 2 (w=1.8): 테스트 항목
    # col 3 (w=1.2): Pass/Fail
    # 합쳐서 표현
    tc_rows = [
        # (col0, col1, col2, col3)
        ('TC-004 / BVA', 'UC-1', '주행 중 해제 차단 (3.0 km/h) ★ASIL B', 'PASS'),
        ('TC-005 / BVA', 'UC-1', '주행 중 해제 허용 (2.9 km/h) ★ASIL B', 'PASS'),
        ('TC-008 / BVA', 'UC-2', '속도 초과 자동잠금 동작 (15.0 km/h)', 'PASS'),
        ('TC-009 / DT',  'UC-3', '충돌 시 비상 해제 — 복합조건 ★ASIL B', 'PASS'),
        ('TC-010 / DT',  'UC-3', '충돌 단일조건 무시 (오동작 방지) ★ASIL B', 'PASS'),
        ('TC-011 / ST',  'UC-4', 'ECU 상태 Init→Normal→Lock 핸들 차단', 'PASS'),
        ('TC-012 / PW',  'UC-5', '후방 객체(거리/속도/시간) 복합 감지', 'PASS'),
        ('TC-006 / FI',  'UC-1', 'ECU 통신 단절 재전송 & Fail-Safe', 'PASS'),
        ('TC-015 / FI',  'UC-6', '점유 센서 에러 시 보호 로직', 'PASS'),
        ('나머지 6건',   'ALL UC', 'EP/BVA/DT 추가 시나리오 (상세: ReqTest.md)', 'PASS'),
    ]

    # shape index 구조 분석:
    # Header rect/textbox: [5~12] (총 8개: 4 rect + 4 textbox)
    # Data rows: [13~52] (총 40개: 5행 × 4col × 2(rect+textbox))
    # row 0: shapes [13~20], row 1: [21~28], row 2: [29~36], row 3: [37~44], row 4: [45~52]
    # 각 row 내: col0_rect, col0_tb, col1_rect, col1_tb, col2_rect, col2_tb, col3_rect, col3_tb

    # 헤더 textbox들의 인덱스: 6, 8, 10, 12
    header_tb_indices = [6, 8, 10, 12]
    shapes = list(slide.shapes)

    # 헤더 교체
    header_labels = ['Test ID / 기법', 'UC', '테스트 항목 (15건 중 핵심)', 'PASS/FAIL']
    for idx, label in zip(header_tb_indices, header_labels):
        replace_shape_text(shapes[idx], label, size=9, bold=True, color=WHITE)

    # 데이터 행 textbox 교체
    # row i: textbox indices = 14, 16, 18, 20 + 8*i
    for row_i in range(5):
        base = 13 + row_i * 8  # rect Col0 index
        tb_col0 = base + 1
        tb_col1 = base + 3
        tb_col2 = base + 5
        tb_col3 = base + 7
        tc = tc_rows[row_i] if row_i < len(tc_rows) else ('', '', '', '')
        is_asil = '★ASIL B' in tc[2]

        # col0: Test ID
        replace_shape_text(shapes[tb_col0], tc[0], size=8.5, bold=is_asil,
                           color=ASIL_B if is_asil else DARK)
        # col1: UC
        replace_shape_text(shapes[tb_col1], tc[1], size=9, bold=False, color=DARK)
        # col2: 항목
        items_text = tc[2].replace(' ★ASIL B', '')
        if is_asil:
            items_text += '  [ASIL B]'
        replace_shape_text(shapes[tb_col2], items_text, size=8.5, bold=False, color=DARK)
        # col3: Pass/Fail
        replace_shape_text(shapes[tb_col3], tc[3], size=10, bold=True, color=PASS_GR)

    # 오른쪽 성과 지표 영역 업데이트 (shapes [54~62])
    # TextBox 55 (정확도) → 단위 테스트
    # TextBox 57 (100%) →  140 건
    # TextBox 58 (처리속도) → 테스트 통과율
    # TextBox 59 (<1ms) → 100% (0 Fail)
    # TextBox 60 (테스트 통과율) → 시스템 TC
    # TextBox 61 (100%) → 15 건
    # TextBox 62 (커버리지) → 커버리지
    # TextBox 63 (100%) → 100% (Branch)
    kpi_map = {
        '정확도': '단위 테스트',
        '100%': '140 건',
        '처리 속도': '통과율',
        '<1ms': '100%  (0 Fail)',
        '테스트 통과율': '시스템 TC',
        '커버리지': '커버리지',
    }
    # 단순하게: 미리 알려진 텍스트를 찾아 교체
    replacements = {
        '정확도': ('단위 테스트', 10, False, BLUE_DARK),
        '처리 속도': ('통과율', 10, False, BLUE_DARK),
        '테스트 통과율': ('시스템 TC', 10, False, BLUE_DARK),
        '커버리지': ('커버리지', 10, False, BLUE_DARK),
    }
    big_replacements = {
        '<1ms': ('100%\n(0 Fail)', 18, True, PASS_GR),
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if t in replacements:
            label, sz, bd, col = replacements[t]
            replace_shape_text(shape, label, size=sz, bold=bd, color=col)
        elif t in big_replacements:
            label, sz, bd, col = big_replacements[t]
            replace_shape_text(shape, label, size=sz, bold=bd, color=col)
        elif t == '100%' and '성과 지표' not in t:
            # 여러 개의 100% 텍스트 box가 있음. 위치로 구분
            l_in = round(shape.left / 914400, 1) if shape.left else 0
            t_in = round(shape.top / 914400, 1) if shape.top else 0
            if t_in < 3.0:  # 위쪽 → 단위 테스트 140건
                replace_shape_text(shape, '140 건', size=22, bold=True, color=BLUE_DARK)
            elif 4.0 < t_in < 5.0:  # 중간 → 15건 시스템 TC
                replace_shape_text(shape, '15 건', size=22, bold=True, color=ASIL_B)
            # 아래 100% → 커버리지 100% 그대로


# ─── SLIDE 8: 학습 내용 & 교훈 ──────────────────────────────────────────────

def fix_slide8_lessons(slide):
    print("  [Slide 8] 학습 내용 보강")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        if 'ISO 26262/ASPICE 산출물' in t:
            replace_shape_multiline(shape, [
                ('① ISO 26262 기반 기능 안전 분석 (HAZOP / FMEA) 및 ASIL B 수준 안전 목표 도출', 10, False, DARK),
                ('② SRS → SwDD → Code → Test 양방향 추적성 확보 (ASPICE SWE.1~SWE.6)', 10, False, DARK),
                ('③ V-모델 기반 설계-구현 워크플로우 — 요구사항 작성 전에 테스트 케이스 먼저 정의 (TDD)', 10, False, DARK),
                ('④ Fail-Safe 상태(SAFE_LOCKED / EMERGENCY_RELEASED) 설계 및 우선순위 정책 수립', 10, False, DARK),
            ])

        elif 'GitHub Action' in t and 'CI/CT' in t:
            replace_shape_multiline(shape, [
                ('① GitHub Actions CI/CD 파이프라인: 빌드 → GTest → Cppcheck → lcov 자동화', 10, False, DARK),
                ('② Lizard 복잡도 분석으로 Cyclomatic Complexity ≤10 준수 자동 검증', 10, False, DARK),
                ('③ GoogleTest Mock 기반 엣지 케이스(Fault Injection, Null 입력) 시뮬레이션', 10, False, DARK),
                ('④ Branch / MC/DC Coverage 100% 달성 — 140개 단위 테스트 케이스 모두 PASS', 10, False, DARK),
            ])

        elif '정상 동작 외의 복잡한 예외 조건' in t:
            replace_shape_multiline(shape, [
                ('▸ 과제 1: 복합 센서 예외 조건 (속도 오류+충돌+핸들 동시 입력) 처리 시 상태 충돌', 10, False, DARK),
                ('▸ 과제 2: 단위 테스트로 검증 어려운 시스템 수준 통합 시나리오 구축', 10, False, DARK),
                ('▸ 과제 3: Door ECU ACK 미수신 상황의 재전송/Fail-Safe 로직 복잡성 증가', 10, False, DARK),
            ])

        elif '제어 우선순위를 반영한 Decision Table' in t:
            replace_shape_multiline(shape, [
                ('▸ 해결 1: CrashSignal > SensorFault > DriverInput > AutoLock 우선순위 Decision Table', 10, False, DARK),
                ('▸ 해결 2: GTest Mock으로 ECU 미응답/센서 고장 시나리오 Fault Injection 자동화', 10, False, DARK),
                ('▸ 해결 3: 상태 전이 다이어그램(State Diagram)을 코드 구조에 1:1 대응시켜 가독성 확보', 10, False, DARK),
            ])


# ─── SLIDE 9: 향후 계획 ───────────────────────────────────────────────────────

def fix_slide9_future(slide):
    print("  [Slide 9] 향후 계획 보강")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        if 'CAN 통신 프로토콜 통합 제어' in t and '단계' not in t:
            replace_shape_text(shape,
                'CAN / CAN-FD 기반 도어 ECU 실시간 제어\n'
                '다중 제어기 통합 시그널 최적화 연동\n'
                '★ 차량 네트워크 통신 스택 통합 검증',
                size=10)

        elif '동적 주행 센서' in t and '단계' not in t:
            replace_shape_text(shape,
                'LiDAR / Radar 복합 센서 데이터 퓨전\n'
                'TTC(Time-to-Collision) 기반 위험도 고도화\n'
                '★ 오탐률 최소화 히스테리시스 알고리즘 고도화',
                size=10)

        elif 'T-05 도어 제어 ECU' in t and '단계' not in t:
            replace_shape_text(shape,
                '도어 제어 ECU 타겟 하드웨어 포팅 (AUTOSAR RTE)\n'
                'HIL(Hardware-in-the-Loop) 환경 검증\n'
                '★ ISO 26262 ASIL B 인증 준비',
                size=10)

        elif '현대모비스 차세대' in t:
            replace_shape_multiline(shape, [
                ('🚗 현대모비스 차세대 통합 하차 안전 제어기 모듈 적용 가능', 11, True, BLUE_DARK),
                ('• 기존 기계식 차일드 락 대체 → 비용 절감 + 안전성 향상', 10, False, DARK),
                ('• ASIL B 인증으로 글로벌 완성차 OEM 납품 가능', 10, False, DARK),
                ('• V2X 연동 확장 시 원격 잠금/해제 상업화 가능', 10, False, DARK),
            ])

        elif 'BLE 기반' in t:
            replace_shape_multiline(shape, [
                ('💡 아이디어 1: BLE 스마트폰 앱 연동 강제 오버라이드', 11, True, BLUE_MID),
                ('   • 부모 앱에서 원격 잠금/해제 + 긴급 해제 알림', 10, False, DARK),
                ('💡 아이디어 2: AI 기반 영유아 방치 감지', 11, True, BLUE_MID),
                ('   • 좌석 점유 + 온도 + 시간 복합 분석 → 방치 사고 예방', 10, False, DARK),
                ('💡 아이디어 3: OTA 업데이트로 임계값 동적 조정', 11, True, BLUE_MID),
                ('   • Safety Goal 변경 없이 파라미터($UNLOCK_INHIBIT_SPEED 등) 업데이트', 10, False, DARK),
            ])


# ─── SLIDE 2: 팀원 카드 ──────────────────────────────────────────────────────

def fix_slide2_team(slide):
    """기존 팀원 내용 텍스트 교체 (shape 찾아 값 변경)"""
    print("  [Slide 2] 팀원 역할 교체")

    # shape 이름으로 교체하기 어려우므로 텍스트 패턴으로 찾기
    member_map = {
        '박기준': (
            '박기준 (팀장)',
            '프로젝트 PM & Safety 모듈\n(PKJthecreator)',
            'F-07 RearRiskEvaluation\nF-08 RearRiskProtectionController\n'
            '• 후방 위험 평가/보호 ASIL B 로직 구현\n'
            '• PM 총괄 (PR #15, #16)'
        ),
        '김도균': (
            '김도균',
            'HMI & 좌석 알림 개발자\n(codingSkeleton6478)',
            'F-06 HmiAndEventLogger\nF-09 RearSeatOccupancyAlert\n'
            '• 경고/알림 HMI 출력 로직 구현\n'
            '• 이벤트 로거 구현 (PR #24)'
        ),
        '김이안': (
            '김이안',
            '상태 결정 로직 개발자\n(iank1m)',
            'F-02 ChildLockStateDecision\n'
            '• 핵심 차일드락 상태 전이 로직 단독 구현\n'
            '• ISO 26262 ASIL 설계 반영 (PR #14)'
        ),
        '이한결': (
            '이한결',
            '프로젝트 메인테이너 & F-01\n(hangyeoli)',
            'F-01 InputMonitorAndValidator\n'
            '• 전체 Issues 생성/관리 (#4~#13)\n'
            '• CMake 빌드 구성 (PR #17)\n'
            '• Scope Fix & CTest (PR #21, #26)'
        ),
        '이승욱': (
            '이승욱',
            '상태 유지 & 알림 + UML 담당\n(josephuk77)',
            'F-05 StatePersistenceManager\nF-10 IgnitionOffStatusAlert\n'
            '• 상태 저장/복구 및 시동 OFF 알림 (PR #19, #20)\n'
            '• UML 다이어그램 문서화 (PR #2)'
        ),
        '박찬석': (
            '박찬석',
            'TDD 기반 제어 로직 개발자\n(p-chanseok)',
            'F-03 DoorEcuCommandHandler\nF-04 RearDoorOpenBlockHandler\n'
            '• TDD 방식 Door ECU 제어 구현\n'
            '• 69/69 테스트 케이스 통과 (PR #22, #23)'
        ),
    }

    # 슬라이드의 textbox들에서 이름 패턴 찾아 교체
    name_shapes = {}  # '이름' → [shape, shape, ...]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        for key in member_map:
            if t == key or t == f'{key} (팀장)':
                name_shapes[key] = shape
                break

    for key, (name, role, detail) in member_map.items():
        if key not in name_shapes:
            continue
        shape = name_shapes[key]
        # 이름 shape를 이름+역할로 교체
        replace_shape_text(shape, name, size=13, bold=True, color=WHITE)

    # 역할/업무 텍스트들도 교체 (기존 "프로젝트 리더..." 패턴 제거)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if '프로젝트 리더 & 아키텍트' in t:
            # 어느 card인지 위치로 추정 불가 → 모두 교체 필요
            replace_shape_text(shape, '', size=9)
        elif '도어 및 하차 안전 제어 로직 기획/구현' in t:
            replace_shape_text(shape, '', size=9)


# ─── 팀명 교체 ────────────────────────────────────────────────────────────────

def fix_slide10(slide):
    print("  [Slide 10] 팀명 교체")
    for shape in slide.shapes:
        if shape.has_text_frame and '[팀명]' in shape.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if '[팀명]' in run.text:
                        run.text = run.text.replace('[팀명]', '전차 팀')


# ─── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    print(f'Loading  : {SRC}')
    prs = Presentation(SRC)

    fix_slide2_team(prs.slides[1])
    fix_slide3_problem(prs.slides[2])
    fix_slide5_arch(prs.slides[4])
    fix_slide6_impl(prs.slides[5])
    fix_slide7_test(prs.slides[6])
    fix_slide8_lessons(prs.slides[7])
    fix_slide9_future(prs.slides[8])
    fix_slide10(prs.slides[9])

    print(f'Saving   : {OUT}')
    prs.save(OUT)
    print('✅ Done!')


if __name__ == '__main__':
    main()
