"""Analyze shapes in key slides and check PIL availability"""
from pptx import Presentation
from pptx.util import Inches
import os

# Check PIL
try:
    from PIL import Image
    print("PIL: AVAILABLE")
except ImportError:
    print("PIL: NOT AVAILABLE")

SRC = r'c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\docs\MBP-T05_ChildLock_결과발표_Updated.pptx'
prs = Presentation(SRC)

SLIDES_TO_ANALYZE = [2, 4, 5, 6, 7, 8]  # 0-indexed (slides 3,5,6,7,8,9 in 1-indexed)

for si in SLIDES_TO_ANALYZE:
    slide = prs.slides[si]
    print(f"\n=== Slide {si+1} (layout={slide.slide_layout.name}) ===")
    for j, shape in enumerate(slide.shapes):
        stype = shape.shape_type
        l = round(shape.left/914400, 2) if shape.left else 0
        t = round(shape.top/914400, 2) if shape.top else 0
        w = round(shape.width/914400, 2) if shape.width else 0
        h = round(shape.height/914400, 2) if shape.height else 0
        txt = ''
        if shape.has_text_frame:
            txt = shape.text.strip().replace('\n',' ')[:80]
        if stype == 19:  # TABLE
            tbl = shape.table
            print(f"  [{j}] TABLE ({tbl.rows.__len__()} rows x {tbl.columns.__len__()} cols) pos=({l},{t}) size={w}x{h}")
            for ri, row in enumerate(tbl.rows):
                row_texts = [cell.text.strip().replace('\n',' ')[:30] for cell in row.cells]
                print(f"       row[{ri}]: {row_texts}")
        else:
            print(f"  [{j}] type={stype} name={shape.name!r} ({l},{t}) {w}x{h}  text={txt!r}")
