from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation(r'docs\MBP-T05_ChildLock_결과발표_Updated.pptx')
sw = prs.slide_width
sh = prs.slide_height
print(f'Slide size: {sw.inches:.2f}" x {sh.inches:.2f}"  ({sw} x {sh} EMU)')
print()
for si, slide in enumerate(prs.slides):
    print(f'=== Slide {si+1} (layout: {slide.slide_layout.name}) ===')
    for j, shape in enumerate(slide.shapes):
        txt = ''
        if shape.has_text_frame:
            txt = shape.text.strip().replace('\n',' ')[:70]
        l_in = shape.left / 914400 if shape.left else 0
        t_in = shape.top / 914400 if shape.top else 0
        w_in = shape.width / 914400 if shape.width else 0
        h_in = shape.height / 914400 if shape.height else 0
        print(f'  [{j}] type={shape.shape_type} name={shape.name!r} ({l_in:.2f},{t_in:.2f})  {w_in:.2f}x{h_in:.2f}  text={txt!r}')
    print()
