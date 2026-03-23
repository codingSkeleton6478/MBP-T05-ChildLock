import collections
import collections.abc
import pptx
collections.Container = collections.abc.Container
collections.Mapping = collections.abc.Mapping
collections.MutableMapping = collections.abc.MutableMapping
collections.Iterable = collections.abc.Iterable
collections.MutableSet = collections.abc.MutableSet
collections.Callable = collections.abc.Callable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs_path = r"c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\docs\주제1_5팀_PBL최종_제출자료.pptx"
prs = Presentation(prs_path)
slide = prs.slides[5]

# 1. Keep only titles and background shapes
shapes_to_keep_names = [
    "Rectangle 1",
    "TextBox 2",
    "Rectangle 3",
    "Rounded Rectangle 5",
    "Rounded Rectangle 15",
    "TextBox 21",
    "TextBox 30"
]

shapes_to_delete = []
for shape in slide.shapes:
    if shape.name not in shapes_to_keep_names:
        shapes_to_delete.append(shape)

for shape in shapes_to_delete:
    element = shape.element
    element.getparent().remove(element)

# Left Column scaling: 0.8
scale = 0.8

# 2. Left Column: Architecture Diagram
arch_y = 2.0
box_h = 0.8 * scale
arch_w = 1.8 * scale
gap = 0.5 * scale

# X positions for Arch boxes
x1 = 0.5
x2 = x1 + arch_w + gap
x3 = x2 + arch_w + 0.2 + gap

# Sensor Block
b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x1), Inches(arch_y), Inches(arch_w), Inches(box_h))
b1.fill.solid()
b1.fill.fore_color.rgb = RGBColor(220, 230, 242)
b1.line.color.rgb = RGBColor(79, 129, 189)
tf = b1.text_frame
tf.text = "[인지]\n데이터 입력"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[1].font.size = Pt(10)
tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[1].alignment = PP_ALIGN.CENTER

slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x1 + arch_w + 0.05), Inches(arch_y + box_h/2 - 0.1), Inches(gap - 0.1), Inches(0.2))

# ECU Block
b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x2), Inches(arch_y), Inches(arch_w + 0.2), Inches(box_h))
b2.fill.solid()
b2.fill.fore_color.rgb = RGBColor(253, 233, 217)
b2.line.color.rgb = RGBColor(247, 150, 70)
tf = b2.text_frame
tf.text = "[판단]\nF-07/F-08 Logic"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[1].font.size = Pt(10)
tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[1].alignment = PP_ALIGN.CENTER

slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x2 + arch_w + 0.2 + 0.05), Inches(arch_y + box_h/2 - 0.1), Inches(gap - 0.1), Inches(0.2))

# Actuator Block
b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x3), Inches(arch_y), Inches(arch_w), Inches(box_h))
b3.fill.solid()
b3.fill.fore_color.rgb = RGBColor(235, 241, 222)
b3.line.color.rgb = RGBColor(155, 187, 89)
tf = b3.text_frame
tf.text = "[제어]\nChild Lock Motor"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[1].font.size = Pt(10)
tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[1].alignment = PP_ALIGN.CENTER

# 3. Left Column: Flowchart
flow_y_start = 2.9 # Moved up
center_x = 3.2 # Flowchart center X

# Start Block (정차 중)
flow_start = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - (0.6*scale)), Inches(flow_y_start), Inches(1.2*scale), Inches(0.5*scale))
flow_start.fill.solid()
flow_start.fill.fore_color.rgb = RGBColor(200, 200, 200)
tf = flow_start.text_frame
tf.text = "정차 중"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Down arrow
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(center_x - 0.1), Inches(flow_y_start + 0.5*scale), Inches(0.2), Inches(0.35))

# Decision 1: Sensor Valid?
d1_y = flow_y_start + 0.5*scale + 0.35
d1_w = 2.4 * scale
d1_h = 1.0 * scale
dec_1 = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(center_x - d1_w/2), Inches(d1_y), Inches(d1_w), Inches(d1_h))
dec_1.fill.solid()
dec_1.fill.fore_color.rgb = RGBColor(255, 255, 204)
tf = dec_1.text_frame
tf.text = "센서 정상\n(riskValid)"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# No arrow -> SAFE_LOCKED
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(center_x + d1_w/2), Inches(d1_y + d1_h/2 - 0.1), Inches(0.5), Inches(0.2))
tb_no = slide.shapes.add_textbox(Inches(center_x + d1_w/2), Inches(d1_y + d1_h/2 - 0.3), Inches(0.5), Inches(0.3))
tb_no.text_frame.text = "No"
tb_no.text_frame.paragraphs[0].font.size = Pt(10)

safe_w = 2.0 * scale
safe_h = 0.6 * scale
safe_state = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(center_x + d1_w/2 + 0.6), Inches(d1_y + d1_h/2 - safe_h/2), Inches(safe_w), Inches(safe_h))
safe_state.fill.solid()
safe_state.fill.fore_color.rgb = RGBColor(255, 153, 153)
tf = safe_state.text_frame
tf.text = "SAFE_LOCKED\n(강제 잠금 + 경고)"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Yes arrow v
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(center_x - 0.1), Inches(d1_y + d1_h), Inches(0.2), Inches(0.4))
tb_yes = slide.shapes.add_textbox(Inches(center_x + 0.1), Inches(d1_y + d1_h), Inches(0.5), Inches(0.3))
tb_yes.text_frame.text = "Yes"
tb_yes.text_frame.paragraphs[0].font.size = Pt(10)

# Decision 2: Threat?
d2_y = d1_y + d1_h + 0.4
d2_w = 3.6 * scale
d2_h = 1.1 * scale
dec_2 = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(center_x - d2_w/2), Inches(d2_y), Inches(d2_w), Inches(d2_h))
dec_2.fill.solid()
dec_2.fill.fore_color.rgb = RGBColor(255, 229, 204)
tf = dec_2.text_frame
tf.text = "위험 감지?\nDist <= THRESH\nSpd >= THRESH"
for p in tf.paragraphs:
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

# Yes arrow 2 -> Risk High
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(center_x + d2_w/2), Inches(d2_y + d2_h/2 - 0.1), Inches(0.5), Inches(0.2))
tb_yes2 = slide.shapes.add_textbox(Inches(center_x + d2_w/2), Inches(d2_y + d2_h/2 - 0.3), Inches(0.5), Inches(0.3))
tb_yes2.text_frame.text = "Yes"
tb_yes2.text_frame.paragraphs[0].font.size = Pt(10)

risk_w = 2.0 * scale
risk_h = 0.6 * scale
risk_high = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(center_x + d2_w/2 + 0.6), Inches(d2_y + d2_h/2 - risk_h/2), Inches(risk_w), Inches(risk_h))
risk_high.fill.solid()
risk_high.fill.fore_color.rgb = RGBColor(255, 102, 102) # Red
risk_high.line.color.rgb = RGBColor(200, 50, 50)
tf = risk_high.text_frame
tf.text = "위험 (Risk High)\n잠금 유지+경고"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# No arrow 2 -> Safe (Green Box)
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(center_x - 0.1), Inches(d2_y + d2_h), Inches(0.2), Inches(0.4))
tb_no2 = slide.shapes.add_textbox(Inches(center_x + 0.1), Inches(d2_y + d2_h), Inches(0.5), Inches(0.3))
tb_no2.text_frame.text = "No"
tb_no2.text_frame.paragraphs[0].font.size = Pt(10)

safe_ok_y = d2_y + d2_h + 0.4
safe_ok_w = 2.4 * scale
safe_ok_h = 0.7 * scale
safe_ok = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(center_x - safe_ok_w/2), Inches(safe_ok_y), Inches(safe_ok_w), Inches(safe_ok_h))
safe_ok.fill.solid()
safe_ok.fill.fore_color.rgb = RGBColor(220, 245, 220) # Pastel green
safe_ok.line.color.rgb = RGBColor(0, 100, 0) # Dark green
tf = safe_ok.text_frame
tf.text = "안전 (Safe)\n잠금 해제 허용 (하차 가능)"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0) # Dark green text
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 4. Right Column: Code Snippet (Exact same as before)
code_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(2.2), Inches(5.0), Inches(4.7))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
code_bg.line.color.rgb = RGBColor(80, 80, 80)

tf = code_bg.text_frame
tf.word_wrap = False
tf.clear()

lines = [
    [("/* F-07: Rear Risk Evaluation */", "comment")],
    [("bool", "keyword"), (" isRawThreatActive(const RearRiskInput_t *in){", "normal")],
    [("    const ", "keyword"), ("bool", "keyword"), (" isClose = (in->dist <= THRESH);", "normal")],
    [("    const ", "keyword"), ("bool", "keyword"), (" isFast = (in->spd >= THRESH);", "normal")],
    [("    return", "keyword"), (" (isClose && isFast);", "normal")],
    [("}", "normal")],
    [("", "normal")],
    [("/* F-08: Protection Controller */", "comment")],
    [("void", "keyword"), (" F08_Run(const Input_t* in, Output_t* out){", "normal")],
    [("    if", "keyword"), (" (!in->riskValid) { ", "normal"), ("// 센서 오류", "comment")],
    [("        out->targetCLState = CL_STATE_ON;", "normal")],
    [("        out->warningSound = ", "normal"), ("true", "keyword"), (";", "normal")],
    [("    }", "normal")],
    [("    else if", "keyword"), (" (in->riskHigh) { ", "normal"), ("// 위험 단계", "comment")],
    [("        out->targetCLState = CL_STATE_ON;", "normal")],
    [("        out->warningSound = ", "normal"), ("true", "keyword"), (";", "normal")],
    [("    }", "normal")],
    [("}", "normal")]
]

colors = {
    "keyword": RGBColor(86, 156, 214),
    "comment": RGBColor(87, 166, 74),
    "normal": RGBColor(212, 212, 212)
}

for i, line_runs in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(2)
    p.space_before = Pt(0)
    for text, style in line_runs:
        run = p.add_run()
        run.text = text
        run.font.name = "Consolas"
        run.font.size = Pt(11)
        run.font.color.rgb = colors.get(style, colors["normal"])


# 5. Connectors/Leader Lines mapping diagram to code
arrow1 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.4), Inches(2.9), Inches(0.6), Inches(0.3))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = RGBColor(200, 200, 200)
arrow1.line.color.rgb = RGBColor(200, 200, 200)

arrow2 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.4), Inches(4.3), Inches(0.6), Inches(0.3))
arrow2.line.color.rgb = RGBColor(255, 153, 153)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = RGBColor(255, 153, 153)

arrow3 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.4), Inches(5.3), Inches(0.6), Inches(0.3))
arrow3.line.color.rgb = RGBColor(255, 102, 102)
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = RGBColor(255, 102, 102)

prs.save(prs_path)
print("PPT updated successfully.")
