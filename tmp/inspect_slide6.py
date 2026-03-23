import collections
import collections.abc
import pptx
collections.Container = collections.abc.Container
collections.Mapping = collections.abc.Mapping
collections.MutableMapping = collections.abc.MutableMapping
collections.Iterable = collections.abc.Iterable
collections.MutableSet = collections.abc.MutableSet
collections.Callable = collections.abc.Callable

from pptx import Presentation

prs_path = r"c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\docs\주제1_5팀_PBL최종_제출자료.pptx"
prs = Presentation(prs_path)
slide = prs.slides[5]

with open(r"c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\tmp\shapes.txt", "w", encoding="utf-8") as f:
    for i, shape in enumerate(slide.shapes):
        f.write(f"Index: {i}\n")
        f.write(f"  Shape ID: {shape.shape_id}\n")
        f.write(f"  Name: {shape.name}\n")
        f.write(f"  Type: {shape.shape_type}\n")
        if hasattr(shape, "text"):
            f.write(f"  Text: {shape.text[:50]}\n")
        f.write("---\n")
