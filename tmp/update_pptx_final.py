"""
MBP-T05 Child Lock - PPT 전면 개선 스크립트 (학회 발표 스타일)
Output: docs/MBP-T05_ChildLock_결과발표_Final.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── 색상 팔레트 ─────────────────────────────────────────────────────────────
def rgb(r, g, b):
    return RGBColor(r, g, b)

BLUE_DARK  = rgb(0x00, 0x35, 0x7e)
BLUE_MID   = rgb(0x00, 0x5b, 0xb5)
BLUE_LIGHT = rgb(0xe8, 0xf0, 0xfe)
ACCENT     = rgb(0x00, 0xa8, 0xe8)
WHITE      = rgb(0xff, 0xff, 0xff)
DARK_TEXT  = rgb(0x1a, 0x1a, 0x2e)
GRAY_BG    = rgb(0xf4, 0xf6, 0xf9)
GREEN      = rgb(0x00, 0x96, 0x60)
ASIL_B     = rgb(0xff, 0x6b, 0x35)
LITE_BLUE  = rgb(0xcc, 0xe5, 0xff)

def hex6(color):
    """RGBColor to 6-char hex string"""
    return '%02X%02X%02X' % (color[0], color[1], color[2])

BASE = r'c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock'
DIAG = os.path.join(BASE, 'docs', 'diagram')
SRC  = os.path.join(BASE, 'docs', 'MBP-T05_ChildLock_결과발표_Updated.pptx')
OUT  = os.path.join(BASE, 'docs', 'MBP-T05_ChildLock_결과발표_Final.pptx')

# ─── 유틸리티 ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    """단색 사각형 배경 추가"""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 선 없음
    return shape


def add_tb(slide, left, top, width, height, text,
           size=11, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    """텍스트 박스 추가 (배경 투명)"""
    if color is None:
        color = DARK_TEXT
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Malgun Gothic'
    return txb


def add_tb_multiline(slide, left, top, width, height, lines,
                     size=11, bold=False, color=None, align=PP_ALIGN.LEFT):
    """여러 줄 텍스트 박스"""
    if color is None:
        color = DARK_TEXT
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (line_text, line_bold, line_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size if line_size else size)
        run.font.bold = line_bold if line_bold is not None else bold
        run.font.color.rgb = color
        run.font.name = 'Malgun Gothic'
    return txb


def fit_image(slide, img_path, left, top, max_w, max_h):
    """이미지를 비율 유지하며 영역 내 중앙 배치"""
    try:
        from PIL import Image
        img = Image.open(img_path)
        iw, ih = img.size
        ratio = min(max_w / iw, max_h / ih)
        nw = iw * ratio
        nh = ih * ratio
        cx = left + (max_w - nw) / 2
        cy = top + (max_h - nh) / 2
        return slide.shapes.add_picture(
            img_path, Inches(cx), Inches(cy), Inches(nw), Inches(nh))
    except ImportError:
        # PIL 없으면 그냥 width로 맞춤
        try:
            return slide.shapes.add_picture(
                img_path, Inches(left), Inches(top), Inches(max_w))
        except Exception as e:
            print(f'  [WARN] Image {os.path.basename(img_path)}: {e}')
            return None
    except Exception as e:
        print(f'  [WARN] Image {os.path.basename(img_path)}: {e}')
        return None


def remove_shapes_with_text(slide, keywords):
    """특정 키워드를 포함한 shape 제거"""
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text.strip()
            if any(kw in t for kw in keywords):
                to_remove.append(shape)
    for s in to_remove:
        try:
            s.element.getparent().remove(s.element)
        except Exception:
            pass


def add_footer(slide, page_num, total=10):
    add_rect(slide, 0, 6.82, 13.33, 0.32, BLUE_DARK)
    add_tb(slide, 0.1, 6.85, 13.0, 0.28,
           f'MOBIUS BOOTCAMP 1기 | PBL 결과발표 | {page_num}/{total}',
           size=8, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Slide 2: 팀원 카드 재구성 ───────────────────────────────────────────────

def rebuild_slide2(slide):
    print('  Slide 2: 팀원 카드 재구성')
    remove_shapes_with_text(slide, [
        '박기준', '김도균', '김이안', '이한결', '이승욱', '박찬석',
        '프로젝트 리더', '담당 업무', '역할', 'Core Developer',
        '도어 및 하차', '단위 테스트'
    ])

    members = [
        ('이한결', 'hangyeoli', '프로젝트 메인테이너 & F-01',
         'F-01 InputMonitorAndValidator',
         ['전체 Issues 생성 관리 (Issue #4~#13)',
          'CMake 빌드 구성 (PR#17)', 'Scope Fix & CTest (PR#26)'],
         BLUE_DARK),
        ('김이안', 'iank1m', '상태 결정 로직 개발자',
         'F-02 ChildLockStateDecision',
         ['차일드락 핵심 상태 전이 로직',
          '단독 구현 (PR#14)',
          'ISO 26262 ASIL 설계 반영'],
         BLUE_MID),
        ('박찬석', 'p-chanseok', 'TDD 기반 제어 로직 개발자',
         'F-03 DoorEcuCommandHandler\nF-04 RearDoorOpenBlockHandler',
         ['TDD 방식 Door ECU 제어 (PR#23)',
          '69/69 테스트 케이스 통과 (PR#22)'],
         BLUE_DARK),
        ('이승욱', 'josephuk77', '상태 유지 & 알림 개발자',
         'F-05 StatePersistenceManager\nF-10 IgnitionOffStatusAlert',
         ['상태 저장/복구 및 시동 OFF 알림 (PR#19, #20)',
          'UML 다이어그램 문서화 (Issue#1/PR#2)'],
         BLUE_MID),
        ('김도균', 'codingSkeleton6478', 'HMI & 좌석 알림 개발자',
         'F-06 HmiAndEventLogger\nF-09 RearSeatOccupancyAlert',
         ['경고/알림 HMI 출력 로직 (PR#24)',
          '이벤트 로거 구현'],
         BLUE_DARK),
        ('박기준 (팀장)', 'PKJthecreator', 'PM & Safety 모듈 개발자',
         'F-07 RearRiskEvaluation\nF-08 RearRiskProtectionController',
         ['후방 위험 평가/보호 핵심 ASIL B 로직',
          '(PR#15, PR#16) + PM 총괄'],
         BLUE_MID),
    ]

    # 3×2 카드 레이아웃
    card_w, card_h = 4.1, 2.55
    gap = 0.13
    col_x = [0.13, col_x1 := 0.13 + card_w + gap, 0.13 + 2*(card_w+gap)]
    row_y = [1.32, 1.32 + card_h + 0.12]
    col_positions = [0.13, 0.13 + card_w + gap, 0.13 + 2*(card_w+gap)]

    for i, m in enumerate(members):
        col = i % 3
        row = i // 3
        x = col_positions[col]
        y = row_y[row]
        name, gid, role, module, details, hdr_color = m

        # 헤더
        add_rect(slide, x, y, card_w, 0.8, hdr_color)
        add_tb(slide, x+0.08, y+0.06, card_w-0.16, 0.38,
               name, size=13, bold=True, color=WHITE)
        add_tb(slide, x+0.08, y+0.44, card_w-0.16, 0.32,
               f'@{gid}  |  {role}',
               size=8.5, bold=False, color=LITE_BLUE)

        # 본문 배경
        add_rect(slide, x, y+0.8, card_w, card_h-0.8, GRAY_BG)

        # 모듈명
        add_tb(slide, x+0.1, y+0.85, card_w-0.2, 0.45,
               module, size=9, bold=True, color=hdr_color)

        # 상세 내용
        detail_y = y + 1.35
        for line in details:
            add_tb(slide, x+0.12, detail_y, card_w-0.24, 0.28,
                   f'• {line}', size=8, color=DARK_TEXT)
            detail_y += 0.28


# ─── Slide 5: 아키텍처 다이어그램 추가 ──────────────────────────────────────

def rebuild_slide5(slide):
    print('  Slide 5: 아키텍처 다이어그램 삽입')
    remove_shapes_with_text(slide, [
        '시스템 아키텍처', '이미지를 삽입', 'ARCHITECTURE', 'DIAGRAM'
    ])

    # 수직 구분선
    add_rect(slide, 5.85, 1.2, 0.04, 5.5, BLUE_LIGHT)

    # === 오른쪽: 상태 다이어그램 ===
    add_rect(slide, 5.95, 1.2, 3.6, 0.32, BLUE_DARK)
    add_tb(slide, 5.95, 1.2, 3.6, 0.32,
           '🔷 Child Lock 상태 다이어그램',
           size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    state_img = os.path.join(DIAG, 'state', 'childlock_state.png')
    fit_image(slide, state_img, 5.95, 1.52, 3.6, 2.2)

    # === 오른쪽: Use Case 다이어그램 ===
    add_rect(slide, 9.65, 1.2, 3.55, 0.32, BLUE_MID)
    add_tb(slide, 9.65, 1.2, 3.55, 0.32,
           '🔷 유스케이스 다이어그램 (Normal)',
           size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    uc_img = os.path.join(DIAG, 'usecase', 'usecase_normal.png')
    fit_image(slide, uc_img, 9.65, 1.52, 3.55, 2.2)

    # === CI/CT 파이프라인 섹션 ===
    add_rect(slide, 5.95, 3.8, 7.25, 0.32, ACCENT)
    add_tb(slide, 5.95, 3.8, 7.25, 0.32,
           '⚙️  CI/CT 자동화 파이프라인 (GitHub Actions)',
           size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    ci_items = [
        ('Static Analysis', 'Cppcheck', BLUE_DARK),
        ('Complexity', 'Lizard (≤10)', BLUE_MID),
        ('Coverage', 'lcov / gcov', BLUE_DARK),
        ('Code Metrics', 'cloc', BLUE_MID),
        ('Unit Test', 'GTest v1.17.0\n(140 cases)', BLUE_DARK),
        ('Build System', 'CMake + CTest', BLUE_MID),
    ]
    cw2 = 2.36
    gap2 = 0.04
    for idx, (lbl, val, cc) in enumerate(ci_items):
        col2 = idx % 3
        row2 = idx // 3
        cx2 = 5.95 + col2 * (cw2 + gap2)
        cy2 = 4.17 + row2 * 0.6
        add_rect(slide, cx2, cy2, cw2, 0.56, BLUE_LIGHT)
        add_rect(slide, cx2, cy2, 0.04, 0.56, cc)
        add_tb(slide, cx2+0.1, cy2+0.04, cw2-0.14, 0.24,
               lbl, size=8.5, bold=True, color=cc)
        add_tb(slide, cx2+0.1, cy2+0.28, cw2-0.14, 0.24,
               val, size=8.5, color=DARK_TEXT)

    # 품질 기준
    add_rect(slide, 5.95, 5.38, 7.25, 0.3, rgb(0xf0, 0xf4, 0xff))
    add_tb(slide, 6.0, 5.4, 7.15, 0.26,
           '📌 품질 기준: 함수 ≤80 lines  |  Dynamic Memory 금지  |  Defensive Programming  |  MISRA C/C++ 준수',
           size=8, color=BLUE_DARK)


# ─── Slide 6: 구현 결과 이미지 ──────────────────────────────────────────────

def rebuild_slide6(slide):
    print('  Slide 6: 구현 결과 이미지 삽입')
    remove_shapes_with_text(slide, [
        '스크린샷', '화면', '안전 상태', 'AI 기반', 'MOBIUS BOOTCAMP'
    ])

    imgs = [
        (os.path.join(DIAG, 'flow_chart', 'FG-01_door_control.png'),
         'FG-01: 도어 잠금/해제 제어', 'UC-1, 2, 3, 4 반영', BLUE_DARK),
        (os.path.join(DIAG, 'flow_chart', 'FG-02_exit_safety.png'),
         'FG-02: 하차 안전 보조', 'UC-5, 6, 7 반영', BLUE_MID),
        (os.path.join(DIAG, 'sequence', 'sequence_uc1.png'),
         'UC-1 시퀀스: 운전자 ON/OFF 제어', '스위치 → ECU → HMI 흐름', BLUE_DARK),
        (os.path.join(DIAG, 'sequence', 'sequence_uc3.png'),
         'UC-3 시퀀스: 충돌 시 비상 해제', 'ASIL B 안전 해제 경로', ASIL_B),
    ]

    # 2×2 레이아웃
    col_x = [0.1, 6.72]
    row_y = [1.32, 4.1]
    img_w = 6.5
    img_h = 2.6

    for i, (img_path, title, sub, hcolor) in enumerate(imgs):
        col = i % 2
        row = i // 2
        lx = col_x[col]
        ty = row_y[row]

        # 캡션 헤더
        add_rect(slide, lx, ty, img_w, 0.36, hcolor)
        add_tb(slide, lx+0.1, ty+0.04, img_w*0.65, 0.28,
               f'🔷 {title}', size=10, bold=True, color=WHITE)
        add_tb(slide, lx + img_w*0.55, ty+0.06, img_w*0.43, 0.28,
               sub, size=8.5, color=LITE_BLUE, align=PP_ALIGN.RIGHT)

        # 이미지
        fit_image(slide, img_path, lx, ty+0.36, img_w, img_h-0.36)

    # 중앙 구분선
    add_rect(slide, 6.58, 1.32, 0.07, 5.55, BLUE_LIGHT)

    add_footer(slide, 6)


# ─── Slide 7: 테스트 케이스 표 ──────────────────────────────────────────────

def rebuild_slide7(slide):
    print('  Slide 7: 테스트 케이스 표 재구성')
    remove_shapes_with_text(slide, [
        '정상 동작 검증', '테스트 항목', '성과 지표',
        'Pass/Fail', '정확도', '처리 속도', '커버리지', '테스트 통과율',
        'MOBIUS BOOTCAMP'
    ])
    # 테이블 shape 제거
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            to_remove.append(shape)
    for s in to_remove:
        try:
            s.element.getparent().remove(s.element)
        except Exception:
            pass

    # TC 데이터
    tcs = [
        ('TC-001', 'UC-1', '정차 중 수동 스위치 제어 정상 동작', 'EP', 'QM'),
        ('TC-002', 'UC-1', '스위치 채터링 디바운싱 - 49ms (Off-limit)', 'BVA', 'QM'),
        ('TC-003', 'UC-1', '스위치 채터링 디바운싱 - 50ms (On-limit)', 'BVA', 'QM'),
        ('TC-004', 'UC-1', '주행 중 해제 차단 - 3.0km/h (On-limit)', 'BVA', 'ASIL B'),
        ('TC-005', 'UC-1', '주행 중 해제 허용 - 2.9km/h (Off-limit)', 'BVA', 'ASIL B'),
        ('TC-006', 'UC-1', 'ECU 통신 단절 재전송 및 Fail-Safe 안전상태', 'FI', 'QM'),
        ('TC-007', 'UC-2', '속도 초과 자동잠금 미동작 - 14.9km/h', 'BVA', 'QM'),
        ('TC-008', 'UC-2', '속도 초과 자동잠금 동작 - 15.0km/h', 'BVA', 'QM'),
        ('TC-009', 'UC-3', '충돌 시 비상 해제 (복합조건 = True)', 'DT', 'ASIL B'),
        ('TC-010', 'UC-3', '충돌 단일조건 무시 (오동작 방지)', 'DT', 'ASIL B'),
        ('TC-011', 'UC-4', 'ECU 상태 전이(Init→Normal→Lock) 핸들 차단', 'ST', 'QM'),
        ('TC-012', 'UC-5', '후방 객체(거리/속도/시간) 복합 감지', 'PW', 'QM'),
        ('TC-013', 'UC-5', '후방 센서 노이즈/반복 경고 히스테리시스', 'EP', 'QM'),
        ('TC-014', 'UC-6', '승객 점유 및 잠금 OFF 출발 시 알림', 'DT', 'QM'),
        ('TC-015', 'UC-6', '점유 센서 에러 발생 시 보호 로직', 'FI', 'QM'),
    ]
    hdrs = ['Test ID', 'UC', '테스트 항목', '기법', 'ASIL']
    cws  = [0.85, 0.52, 4.05, 0.55, 0.72]
    total_w = sum(cws)
    sx = (13.33 - total_w) / 2
    sy = 1.32
    hh = 0.36
    rh = 0.285

    # 헤더 행
    xp = sx
    for h, cw in zip(hdrs, cws):
        add_rect(slide, xp, sy, cw, hh, BLUE_DARK)
        add_tb(slide, xp+0.03, sy+0.07, cw-0.06, hh-0.1,
               h, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        xp += cw

    # 데이터 행
    for ri, tc in enumerate(tcs):
        ry = sy + hh + ri * rh
        is_asil = tc[4] == 'ASIL B'
        row_bg = BLUE_LIGHT if ri % 2 == 0 else WHITE
        xp = sx
        for ci, (val, cw) in enumerate(zip(tc, cws)):
            cbg = row_bg
            txt_color = DARK_TEXT
            bld = False
            if ci == 4 and is_asil:
                cbg = ASIL_B
                txt_color = WHITE
                bld = True
            add_rect(slide, xp, ry, cw, rh, cbg)
            add_tb(slide, xp+0.03, ry+0.04, cw-0.06, rh-0.06,
                   val, size=8, bold=bld, color=txt_color,
                   align=PP_ALIGN.CENTER)
            xp += cw

    # KPI 박스
    kpi_y = sy + hh + len(tcs) * rh + 0.12
    kpis = [
        ('단위 테스트', '140 건', BLUE_DARK),
        ('통과율', '100%  (0 Fail)', GREEN),
        ('커버리지', '100%  Branch', BLUE_MID),
        ('처리 속도', '< 1 ms', ACCENT),
        ('시스템 TC', '15 건', ASIL_B),
        ('프레임워크', 'GTest v1.17', BLUE_DARK),
    ]
    kw = 2.1
    kg = 0.1
    total_kw = len(kpis)*(kw+kg) - kg
    ksx = (13.33 - total_kw) / 2
    kh = 0.82

    for ki, (lbl, val, kc) in enumerate(kpis):
        kx = ksx + ki*(kw+kg)
        add_rect(slide, kx, kpi_y, kw, 0.28, kc)
        add_tb(slide, kx+0.04, kpi_y+0.04, kw-0.08, 0.22,
               lbl, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, kx, kpi_y+0.28, kw, kh-0.28, rgb(0xf9, 0xfb, 0xff))
        add_tb(slide, kx+0.04, kpi_y+0.3, kw-0.08, kh-0.32,
               val, size=12, bold=True, color=kc, align=PP_ALIGN.CENTER)

    add_footer(slide, 7)


# ─── Slide 10: 팀명 교체 ─────────────────────────────────────────────────────

def fix_slide10(slide):
    print('  Slide 10: 팀명 교체')
    for shape in slide.shapes:
        if shape.has_text_frame and '[팀명]' in shape.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if '[팀명]' in run.text:
                        run.text = run.text.replace('[팀명]', '전차 팀')


# ─── 메인 ────────────────────────────────────────────────────────────────────

def main():
    print(f'Loading: {SRC}')
    prs = Presentation(SRC)

    rebuild_slide2(prs.slides[1])
    rebuild_slide5(prs.slides[4])
    rebuild_slide6(prs.slides[5])
    rebuild_slide7(prs.slides[6])
    fix_slide10(prs.slides[9])

    print(f'Saving: {OUT}')
    prs.save(OUT)
    print('✅ Done!')


if __name__ == '__main__':
    main()
