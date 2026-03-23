import collections
import collections.abc

# Fix for pptx collections.abc deprecation
import pptx
collections.Container = collections.abc.Container
collections.Mapping = collections.abc.Mapping
collections.MutableMapping = collections.abc.MutableMapping
collections.Iterable = collections.abc.Iterable
collections.MutableSet = collections.abc.MutableSet
collections.Callable = collections.abc.Callable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs_path = r"c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\docs\주제1_5팀_PBL최종_제출자료.pptx"
prs = Presentation(prs_path)
slide = prs.slides[5]  # 6th slide

# Clear existing shapes except titles
shapes_to_delete = []
for idx, shape in enumerate(slide.shapes):
    # Keep placeholders (usually main title is placeholder 0)
    if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 0:
        continue
    # Keep anything above 1.3 inches just in case it's a top banner
    if hasattr(shape, "top") and shape.top < Inches(1.3):
        continue
    shapes_to_delete.append(shape)

for shape in shapes_to_delete:
    element = shape.element
    element.getparent().remove(element)

# Left Column Title
left_subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5))
tf = left_subtitle.text_frame
tf.text = "아키텍처 (Architecture)"
p = tf.paragraphs[0]
p.font.bold = True
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(0, 51, 102)

# Block 1: Sensor
b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.0), Inches(1.0))
b1.fill.solid()
b1.fill.fore_color.rgb = RGBColor(220, 230, 242)
b1.line.color.rgb = RGBColor(79, 129, 189)
tf = b1.text_frame
tf.text = "데이터 입력 (Sensor)\nRadar / Ultrasonic"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
if len(tf.paragraphs) > 1:
    tf.paragraphs[1].font.size = Pt(12)
    tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)

# Traceability 1
tb1 = slide.shapes.add_textbox(Inches(0.4), Inches(3.6), Inches(2.2), Inches(1.0))
tb1.text_frame.word_wrap = True
p = tb1.text_frame.add_paragraph()
p.text = "• ISO 26262 ASIL B\n• HW Fault Detection"
p.font.size = Pt(12)

# Arrow 1
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.6), Inches(2.8), Inches(0.6), Inches(0.4))

# Block 2: ECU
b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), Inches(2.5), Inches(2.0), Inches(1.0))
b2.fill.solid()
b2.fill.fore_color.rgb = RGBColor(253, 233, 217)
b2.line.color.rgb = RGBColor(247, 150, 70)
tf = b2.text_frame
tf.text = "제어 판단 (ECU)\nF-07/F-08 Logic"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
if len(tf.paragraphs) > 1:
    tf.paragraphs[1].font.size = Pt(12)
    tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)

# Traceability 2
tb2 = slide.shapes.add_textbox(Inches(3.2), Inches(3.6), Inches(2.2), Inches(1.0))
tb2.text_frame.word_wrap = True
p = tb2.text_frame.add_paragraph()
p.text = "• A-SPICE SWE.2\n• Risk Eval & Protection"
p.font.size = Pt(12)

# Arrow 2
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.4), Inches(2.8), Inches(0.6), Inches(0.4))

# Block 3: Actuator
b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(2.5), Inches(2.0), Inches(1.0))
b3.fill.solid()
b3.fill.fore_color.rgb = RGBColor(235, 241, 222)
b3.line.color.rgb = RGBColor(155, 187, 89)
tf = b3.text_frame
tf.text = "최종 출력 (Actuator)\nChild Lock Motor"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
if len(tf.paragraphs) > 1:
    tf.paragraphs[1].font.size = Pt(12)
    tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)

# Traceability 3
tb3 = slide.shapes.add_textbox(Inches(6.0), Inches(3.6), Inches(2.2), Inches(1.0))
tb3.text_frame.word_wrap = True
p = tb3.text_frame.add_paragraph()
p.text = "• Safe State Policy\n• SAFE_LOCKED (ON)"
p.font.size = Pt(12)


# Right Column Title
right_subtitle = slide.shapes.add_textbox(Inches(8.4), Inches(1.5), Inches(4.5), Inches(0.5))
tf = right_subtitle.text_frame
tf.text = "핵심 로직 (Core Logic)"
p = tf.paragraphs[0]
p.font.bold = True
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(0, 51, 102)

# Code Snippet Background
code_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.4), Inches(2.2), Inches(4.7), Inches(4.2))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
code_bg.line.color.rgb = RGBColor(80, 80, 80)

tf = code_bg.text_frame
tf.word_wrap = False
tf.clear()

lines = [
    [("/* F-07: Rear Risk Evaluation */", "comment")],
    [("bool", "keyword"), (" isRawThreatActive(const RiskIn_t *in){", "normal")],
    [("    const ", "keyword"), ("bool", "keyword"), (" isClose = (in->dist <= THRESH);", "normal")],
    [("    const ", "keyword"), ("bool", "keyword"), (" isFast = (in->spd >= THRESH);", "normal")],
    [("    return", "keyword"), (" (isClose && isFast);", "normal")],
    [("}", "normal")],
    [("", "normal")],
    [("/* F-08: Protection Controller */", "comment")],
    [("void", "keyword"), (" Ctrl_Run(const In_t* in, Out_t* out){", "normal")],
    [("    if", "keyword"), (" (!in->riskValid) { ", "normal"), ("// Sensor Fault", "comment")],
    [("        out->targetCLState = CL_STATE_ON;", "normal")],
    [("        out->warningSound = ", "normal"), ("true", "keyword"), (";", "normal")],
    [("    }", "normal")],
    [("    else if", "keyword"), (" (in->riskHigh) { ", "normal"), ("// Danger Detected", "comment")],
    [("        out->targetCLState = CL_STATE_ON;", "normal")],
    [("        out->warningMsgId = WARNING_RISK_HIGH;", "normal")],
    [("        out->warningSound = ", "normal"), ("true", "keyword"), (";", "normal")],
    [("    }", "normal")],
    [("}", "normal")]
]

colors = {
    "keyword": RGBColor(86, 156, 214),
    "comment": RGBColor(87, 166, 74),
    "normal": RGBColor(212, 212, 212)
}

for i, line_runs in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    # Adjust spacing to fit
    p.space_after = Pt(2)
    p.space_before = Pt(0)
    for text, style in line_runs:
        run = p.add_run()
        run.text = text
        run.font.name = "Consolas"
        run.font.size = Pt(11)
        run.font.color.rgb = colors.get(style, colors["normal"])

prs.save(prs_path)
print("PPT updated successfully.")
