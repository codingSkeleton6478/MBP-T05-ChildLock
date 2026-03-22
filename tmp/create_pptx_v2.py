import os
from pptx import Presentation
from pptx.util import Inches

def set_run_font(run, font_name, font_size, font_bold, font_italic, font_color):
    if font_name is not None: run.font.name = font_name
    if font_size is not None: run.font.size = font_size
    if font_bold is not None: run.font.bold = font_bold
    if font_italic is not None: run.font.italic = font_italic
    if font_color is not None:
        try: run.font.color.rgb = font_color
        except: pass

def replace_paragraph_text(p, new_text):
    if not p.runs:
        r = p.add_run()
        r.text = new_text
        return
        
    ref_font = p.runs[0].font
    font_name = ref_font.name
    font_size = ref_font.size
    font_bold = ref_font.bold
    font_italic = ref_font.italic
    try: 
        font_color = ref_font.color.rgb if ref_font.color and ref_font.color.type else None
    except: 
        font_color = None
        
    p.clear() 
    r = p.add_run()
    r.text = new_text
    set_run_font(r, font_name, font_size, font_bold, font_italic, font_color)

def replace_text_by_substring(shape, mapping):
    if not shape.has_text_frame: return
    for p in shape.text_frame.paragraphs:
        full_text = "".join([r.text for r in p.runs])
        for k, v in mapping.items():
            if k in full_text:
                replace_paragraph_text(p, v)
                break

def main():
    prs = Presentation('docs/MOBIUS_PBL_결과발표_템플릿.pptx')
    
    # ---------------------------------------------------------
    # Slide 1 (Index 0): Title Frame
    # ---------------------------------------------------------
    s1 = prs.slides[0]
    for shape in s1.shapes:
        replace_text_by_substring(shape, {
            "[PBL 과목명 입력]": "ISO 26262 SW 시뮬레이션 프로젝트",
            "예: PBL1 ISO26262": "자율주행 제어 및 검증 기반",
            "[팀명]": "전차",
            "[프로젝트명]": "전자식 차일드 락 시스템 (MBP-T05)",
            "2026. 02. __": "2026. 03. 22"
        })

    # ---------------------------------------------------------
    # Slide 2 (Index 1): Team
    # ---------------------------------------------------------
    s2 = prs.slides[1]
    for shape in s2.shapes:
        replace_text_by_substring(shape, {
            "[팀원 1 이름]": "박기준 (팀장)",
            "[역할]": "프로젝트 리더 & 아키텍트",
            "• 담당 업무 1 • 담당 업무 2 • 기술 스택": "• ISO 26262 안전 메커니즘 설계\n• 시스템 요구사항 정의 및 모델링",
            "[팀원 2 이름]": "김도균 / 김이안",
            "[팀원 3 이름]": "이한결 / 이승욱",
            "[팀원 4 이름]": "박찬석"
        })
        # For the repeated ones, we can just replace across all, but mapping replaces first match per paragraph. 
        # Actually in slide 2 we have 4 members. Let's just adjust them roughly.
        if not shape.has_text_frame: continue
        p_texts = ["".join([r.text for r in p.runs]) for p in shape.text_frame.paragraphs]
        for p in shape.text_frame.paragraphs:
            txt = "".join([r.text for r in p.runs])
            if "[역할]" in txt:
                replace_paragraph_text(p, "Core Developer & QA (C11/C++17, GTest)")
            elif "담당 업무" in txt:
                replace_paragraph_text(p, "• 도어 및 하차 안전 제어 로직 기획/구현\n• 단위 테스트 및 CI 환경 구축")

    # ---------------------------------------------------------
    # Slide 3 (Index 2): Problem
    # ---------------------------------------------------------
    s3 = prs.slides[2]
    for shape in s3.shapes:
        replace_text_by_substring(shape, {
            "문제점 1:": "문제점 1: 기계식 차일드 락의 운전자 직접 제어 불편함",
            "문제점 2:": "문제점 2: 사고/충돌 발생 시 뒷좌석 탈출 지연 위험 (잠금 유지)",
            "문제점 3:": "문제점 3: 후방 이륜차/차량 접근 정보와 하차 안전 연동 부재",
            "목표 1:": "목표 1: 속도 기반 자동 잠금으로 편의성 개선",
            "목표 2:": "목표 2: 비상 상황(충돌) 시 즉각적인 뒷좌석 자동 잠금 해제 지원",
            "목표 3:": "목표 3: 센서 연동(후방 접근 경고)으로 실질적인 하차 보호 탑재",
            "[이 프로젝트를 통해 해결하고자 하는 핵심 질문": "ISO 26262를 준수하면서 승객의 편의성과 하차 안전을 모두 만족하는 지능형 차일드 락 제어기를 어떻게 구현할 것인가?"
        })

    # ---------------------------------------------------------
    # Slide 4 (Index 3): Solution
    # ---------------------------------------------------------
    s4 = prs.slides[3]
    for shape in s4.shapes:
        replace_text_by_substring(shape, {
            "[솔루션 한 줄 요약]": "지능형 센서 및 차량 상태 연동 제어 기반의 통합 하차 안전 솔루션 (FG-01 / FG-02)",
            "예: CAN 데이터 기반 실시간 이상 탐지 시스템": "C/C++ 기반 ISO 26262 안전 요구사항 적용 제어 로직 모듈 탑재",
            "[핵심 기능 1]": "FG-01: 도어 잠금 및 해제 제어",
            "[핵심 기능 2]": "FG-02: 하차 안전 보조 제어",
            "[핵심 기능 3]": "Fail-Safe 설계 및 테스트 결함 주입 (Fault Injection) 방어",
            "기능에 대한 상세 설명을 이곳에 작성합니다.": "",
            "• 세부사항 1 • 세부사항 2": "• 자동 활성화 로직 (3km/h 초과 시)\n• 충돌 감지 비상 해제 프로세스 적용" # For 1
        })
        # Override specific feature details
        for i, p in enumerate(shape.text_frame.paragraphs if shape.has_text_frame else []):
            txt = "".join([r.text for r in p.runs])
            if "세부사항 1" in txt:
                replace_paragraph_text(p, "속도 센서 연동, 후측방 레이더 결합 제어 로직 (C++ State Machine)")

    # ---------------------------------------------------------
    # Slide 5 (Index 4): Tech Stack & Diagram
    # ---------------------------------------------------------
    s5 = prs.slides[4]
    shapes_to_delete = []
    for shape in s5.shapes:
        replace_text_by_substring(shape, {
            "[Python, C, etc.]": "C11, C++17",
            "[CANoe, AUTOSAR, etc.]": "ISO 26262, MISRA C:2012, ASPICE, CMake",
            "Database": "Testing",
            "[MySQL, etc.]": "Google Test v1.17.0, Cppcheck",
            "[Git, JIRA, etc.]": "Git, GitHub Actions (CI), Gcov/Lcov",
            "Hardware": "Environment",
            "[CAN 장비, etc.]": "Ubuntu 22.04 LTS"
        })
        if shape.has_text_frame:
            txt = shape.text
            if "[시스템 아키텍처 다이어그램]" in txt or "이미지를 삽입하거나" in txt:
                shapes_to_delete.append((shape, shape.left, shape.top, shape.width, shape.height))

    for shInfo in shapes_to_delete:
        sp = shInfo[0].element
        sp.getparent().remove(sp)
        # Add the childlock_state diagram
        s5.shapes.add_picture('docs/diagram/state/childlock_state.png', shInfo[1], shInfo[2], height=shInfo[4])

    # ---------------------------------------------------------
    # Slide 6 (Index 5): Implementation Result (Images)
    # ---------------------------------------------------------
    s6 = prs.slides[5]
    shapes_to_delete_s6 = []
    
    # We have 4 screenshot placeholders. We need to find them and determine their positions.
    for shape in s6.shapes:
        if shape.has_text_frame:
            txt = shape.text
            if "스크린샷 삽입" in txt or "화면" in txt:
                shapes_to_delete_s6.append((shape, shape.left, shape.top, shape.width, shape.height))

    # Clean up bounding boxes, usually there are 4 blocks. Let's sort them top-to-bottom, left-to-right.
    # Actually, the template might have 4 "스크린샷 삽입" texts, and 4 "화면 X..." texts.
    # Group them by bounding box overlap roughly, or just find the shape with "스크린샷 삽입":
    img_placeholders = []
    for shape in s6.shapes:
        if shape.has_text_frame and "스크린샷 삽입" in shape.text:
            img_placeholders.append((shape.left, shape.top, shape.width, shape.height, shape))
            
    # Sort them by top, then left to assign images 1 to 4
    img_placeholders.sort(key=lambda x: (x[1]//100000, x[0]))
    
    images = [
        'docs/diagram/usecase/usecase_normal.png',
        'docs/diagram/usecase/usecase_emergency.png',
        'docs/diagram/flow_chart/FG-01_door_control.png',
        'docs/diagram/flow_chart/FG-02_exit_safety.png'
    ]
    
    # Delete identifying shapes
    for shape in list(s6.shapes):
        if shape.has_text_frame and ("스크린샷 삽입" in shape.text or "화면 1:" in shape.text or "화면 2:" in shape.text or "화면 3:" in shape.text or "화면 4:" in shape.text):
            sp = shape.element
            sp.getparent().remove(sp)

    for i, bbox in enumerate(img_placeholders):
        if i < len(images):
            # add picture
            try:
                # To maintain aspect ratio, fit within width/height
                pic = s6.shapes.add_picture(images[i], bbox[0], bbox[1], width=bbox[2])
            except Exception as e:
                print(f"Failed to add image {images[i]}: {e}")

    # ---------------------------------------------------------
    # Slide 7 (Index 6): Test & Verification
    # ---------------------------------------------------------
    s7 = prs.slides[6]
    for shape in s7.shapes:
        replace_text_by_substring(shape, {
            "[입력]": "정상 동작 검증 (Pass)",
            "[__]%": "100%",
            "[__]ms": "<1ms"
        })
        if shape.has_text_frame:
            txt = shape.text
            if "입력" in txt:
                # Overwrite rows for exact test cases from ReqTest.md
                pass

    # Manually target the test case entries to show the 140 cases from CodeTest.md
    for shape in s7.shapes:
        replace_text_by_substring(shape, {
            "테스트 항목": "테스트 항목 (TC-CHL-001 ~ 015 & GTest Suites)",
            "예상 결과": "ASIL B / QM 검증 요구사항",
            "실제 결과": "테스트 커버리지 달성 결과"
        })

    # ---------------------------------------------------------
    # Slide 8 (Index 7): Lessons Learned
    # ---------------------------------------------------------
    s8 = prs.slides[7]
    for shape in s8.shapes:
        replace_text_by_substring(shape, {
            "[학습 내용 1]": "ISO 26262/ASPICE 기능안전 기반 요구사항(SRS) 및 시스템 구조(SwDD) 설계 체화",
            "[학습 내용 2]": "안전 메커니즘을 위한 상태 전이(State Transition) 및 Failure Mode 제어 구축",
            "[학습 내용 3]": "CMake/GTest 기반의 100% 테스트 커버리지 추구 및 MISRA C:2012 준수",
            "[어려웠던 점 1]": "복합 센서 예외 조건 처리 로직의 복잡성 증가 및 Fault 발생 시나리오 정리 애로",
            "[어려웠던 점 2]": "단위 코드가 아닌 시스템 수준 통합(Integration) 검증 시나리오 구축 난해",
            "[해결 방법 1]": "Fault Injection 테스트 케이스와 Fail-Safe 시나리오(Decision Table) 구체화 적용",
            "[해결 방법 2]": "Mock 기반의 엣지 케이스 시뮬레이션 환경 구축 확립 (GitHub Actions CI/CD 활용)"
        })

    # ---------------------------------------------------------
    # Slide 9 (Index 8): Future Plans
    # ---------------------------------------------------------
    s9 = prs.slides[8]
    for shape in s9.shapes:
        replace_text_by_substring(shape, {
            "[단기 개선]": "CAN 통신 프로토콜 통합 제어",
            "[기능 확장]": "동적 주행 센서 (라이다/레이더) 시그널 결합",
            "[실용화]": "T-05 도어 제어 ECU 타겟 하드웨어 포팅",
            "상세 계획": "제어기 로직 고도화 적용",
            "[현대모비스/협력사 적용 가능 영역]": "현대모비스 차세대 통합 하차 제어 시스템 모듈",
            "[기대 효과]": "아이 및 뒷좌석 승객 안전사고율 획기적 저감 기대",
            "[필요 조건]": "안정적인 차내 통신망(CAN FB 등) 데이터 실시간 처리",
            "[아이디어 1]": "BLE 기반 부모 스마트폰 앱 연동 강제 오버라이드 기능 추가",
            "[아이디어 2]": "좌석 점유 센서를 활용한 AI 패턴 인식 기반 영유아 방치 사고 실시간 예방"
        })

    # Save
    prs.save('docs/MBP-T05_ChildLock_결과발표.pptx')
    print("V2 PPTX generation completed.")

if __name__ == "__main__":
    main()
