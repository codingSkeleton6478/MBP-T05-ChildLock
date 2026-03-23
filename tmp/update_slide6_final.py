import collections
import collections.abc
import pptx
collections.Container = collections.abc.Container
collections.Mapping = collections.abc.Mapping
collections.MutableMapping = collections.abc.MutableMapping
collections.Iterable = collections.abc.Iterable
collections.MutableSet = collections.abc.MutableSet
collections.Callable = collections.abc.Callable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs_path = r"c:\GIThub coding\1. F-07,F-08 개발\MBP-T05-ChildLock\docs\주제1_5팀_PBL최종_제출자료.pptx"
prs = Presentation(prs_path)
slide = prs.slides[5]  # Index 5 is the 6th slide

# 1. Keep only titles and background shapes
shapes_to_keep_names = [
    "Rectangle 1",
    "TextBox 2",
    "Rectangle 3",
    "Rounded Rectangle 5",
    "Rounded Rectangle 15",
    "TextBox 21",
    "TextBox 30"
]

shapes_to_delete = []
for shape in slide.shapes:
    if shape.name not in shapes_to_keep_names:
        shapes_to_delete.append(shape)

for shape in shapes_to_delete:
    element = shape.element
    element.getparent().remove(element)

# 2. Left Column: Architecture Diagram
arch_y = 2.2 # Inches
box_h = 0.8
arch_w = 1.8

# Sensor Block
b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(arch_y), Inches(arch_w), Inches(box_h))
b1.fill.solid()
b1.fill.fore_color.rgb = RGBColor(220, 230, 242)
b1.line.color.rgb = RGBColor(79, 129, 189)
tf = b1.text_frame
tf.text = "[인지]\n데이터 입력 (Sensor)"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[1].font.size = Pt(11)
tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[1].alignment = PP_ALIGN.CENTER

slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.4), Inches(arch_y + 0.3), Inches(0.4), Inches(0.2))

# ECU Block
b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.9), Inches(arch_y), Inches(arch_w + 0.2), Inches(box_h))
b2.fill.solid()
b2.fill.fore_color.rgb = RGBColor(253, 233, 217)
b2.line.color.rgb = RGBColor(247, 150, 70)
tf = b2.text_frame
tf.text = "[판단]\nF-07/F-08 Logic"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[1].font.size = Pt(11)
tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[1].alignment = PP_ALIGN.CENTER

slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.0), Inches(arch_y + 0.3), Inches(0.4), Inches(0.2))

# Actuator Block
b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(arch_y), Inches(arch_w), Inches(box_h))
b3.fill.solid()
b3.fill.fore_color.rgb = RGBColor(235, 241, 222)
b3.line.color.rgb = RGBColor(155, 187, 89)
tf = b3.text_frame
tf.text = "[제어]\nChild Lock Motor"
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[1].font.size = Pt(11)
tf.paragraphs[1].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[1].alignment = PP_ALIGN.CENTER


# 3. Left Column: Flowchart
flow_y_start = 3.6

# Start Block
flow_start = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.8), Inches(flow_y_start), Inches(1.2), Inches(0.5))
flow_start.fill.solid()
flow_start.fill.fore_color.rgb = RGBColor(200, 200, 200)
tf = flow_start.text_frame
tf.text = "주행 중"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.3), Inches(flow_y_start + 0.5), Inches(0.2), Inches(0.3))

# Decision 1: Sensor Valid?
dec_1 = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(2.2), Inches(flow_y_start + 0.9), Inches(2.4), Inches(1.0))
dec_1.fill.solid()
dec_1.fill.fore_color.rgb = RGBColor(255, 255, 204)
tf = dec_1.text_frame
tf.text = "센서 정상\n(riskValid)"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# No arrow
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(flow_y_start + 1.25), Inches(0.6), Inches(0.2))
tb_no = slide.shapes.add_textbox(Inches(4.6), Inches(flow_y_start + 1.0), Inches(0.5), Inches(0.3))
tb_no.text_frame.text = "No"

# Safe State Block
safe_state = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(flow_y_start + 1.1), Inches(2.0), Inches(0.6))
safe_state.fill.solid()
safe_state.fill.fore_color.rgb = RGBColor(255, 153, 153)
tf = safe_state.text_frame
tf.text = "SAFE_LOCKED\n(강제 잠금 + 경고)"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Yes arrow
slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.3), Inches(flow_y_start + 1.9), Inches(0.2), Inches(0.4))
tb_yes = slide.shapes.add_textbox(Inches(3.4), Inches(flow_y_start + 1.9), Inches(0.5), Inches(0.3))
tb_yes.text_frame.text = "Yes"

# Decision 2: Threat?
dec_2 = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(1.7), Inches(flow_y_start + 2.4), Inches(3.4), Inches(1.1))
dec_2.fill.solid()
dec_2.fill.fore_color.rgb = RGBColor(255, 229, 204)
tf = dec_2.text_frame
tf.text = "위험 감지?\nDist <= THRESH &&\nSpd >= THRESH"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Yes arrow 2
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.1), Inches(flow_y_start + 2.8), Inches(0.6), Inches(0.2))
tb_yes2 = slide.shapes.add_textbox(Inches(5.1), Inches(flow_y_start + 2.5), Inches(0.5), Inches(0.3))
tb_yes2.text_frame.text = "Yes"

# Risk High Block
risk_high = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(flow_y_start + 2.65), Inches(1.5), Inches(0.6))
risk_high.fill.solid()
risk_high.fill.fore_color.rgb = RGBColor(255, 102, 102)
tf = risk_high.text_frame
tf.text = "위험 (Risk High)\n잠금 유지 + 경고"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER


# 4. Right Column: Code Snippet
code_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(2.2), Inches(5.0), Inches(4.7))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
code_bg.line.color.rgb = RGBColor(80, 80, 80)

tf = code_bg.text_frame
tf.word_wrap = False
tf.clear()

lines = [
    [("/* F-07: Rear Risk Evaluation */", "comment")],
    [("bool", "keyword"), (" isRawThreatActive(const RearRiskInput_t *in){", "normal")],
    [("    const ", "keyword"), ("bool", "keyword"), (" isClose = (in->dist <= THRESH);", "normal")],
    [("    const ", "keyword"), ("bool", "keyword"), (" isFast = (in->spd >= THRESH);", "normal")],
    [("    return", "keyword"), (" (isClose && isFast);", "normal")],
    [("}", "normal")],
    [("", "normal")],
    [("/* F-08: Protection Controller */", "comment")],
    [("void", "keyword"), (" F08_Run(const Input_t* in, Output_t* out){", "normal")],
    [("    if", "keyword"), (" (!in->riskValid) { ", "normal"), ("// 센서 오류", "comment")],
    [("        out->targetCLState = CL_STATE_ON;", "normal")],
    [("        out->warningSound = ", "normal"), ("true", "keyword"), (";", "normal")],
    [("    }", "normal")],
    [("    else if", "keyword"), (" (in->riskHigh) { ", "normal"), ("// 위험 단계", "comment")],
    [("        out->targetCLState = CL_STATE_ON;", "normal")],
    [("        out->warningSound = ", "normal"), ("true", "keyword"), (";", "normal")],
    [("    }", "normal")],
    [("}", "normal")]
]

colors = {
    "keyword": RGBColor(86, 156, 214),
    "comment": RGBColor(87, 166, 74),
    "normal": RGBColor(212, 212, 212)
}

for i, line_runs in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(2)
    p.space_before = Pt(0)
    for text, style in line_runs:
        run = p.add_run()
        run.text = text
        run.font.name = "Consolas"
        run.font.size = Pt(11)
        run.font.color.rgb = colors.get(style, colors["normal"])


# 5. Connectors/Leader Lines mapping diagram to code
# Arrow 1: Dist && Spd in Flowchart -> F-07 code
# Flowchart Box dec_2 is roughly around Y=6.0 (2.4+3.6), Code F-07 is at top
arrow1 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.4), Inches(3.2), Inches(0.6), Inches(0.3))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = RGBColor(255, 204, 0)
arrow1.line.color.rgb = RGBColor(255, 204, 0)

# Arrow 2: Safe State in Flowchart -> F-08 sensor error code
# Safe state box is at Y=4.7 (1.1+3.6). Code is around line 10 (Y=4.2 roughly)
arrow2 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.4), Inches(4.5), Inches(0.6), Inches(0.3))
arrow2.line.color.rgb = RGBColor(255, 153, 153)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = RGBColor(255, 153, 153)

# Arrow 3: Risk High in Flowchart -> F-08 danger code
# Risk high box is at Y=6.25. Code is around line 14 (Y=5.5)
arrow3 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.4), Inches(5.6), Inches(0.6), Inches(0.3))
arrow3.line.color.rgb = RGBColor(255, 102, 102)
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = RGBColor(255, 102, 102)

prs.save(prs_path)
print("PPT updated successfully.")
